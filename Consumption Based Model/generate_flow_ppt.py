#!/usr/bin/env python3
"""
PayGo Architecture: Already Built vs. Needs to Be Built
Flow-diagram style presentation for CIO audience.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
MID_DARK   = RGBColor(0x10, 0x25, 0x3A)
MID_GRAY   = RGBColor(0x1E, 0x30, 0x44)
DK_GRAY    = RGBColor(0x44, 0x55, 0x66)
BLUE       = RGBColor(0x00, 0x72, 0xC6)
TEAL       = RGBColor(0x00, 0xB4, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LT_GRAY    = RGBColor(0xC0, 0xD0, 0xE0)

BUILT      = RGBColor(0x00, 0xA8, 0x55)   # green  — already built
BUILT_BG   = RGBColor(0x04, 0x1E, 0x0E)
EXTEND     = RGBColor(0xFF, 0xCC, 0x00)   # yellow — extend existing
EXTEND_BG  = RGBColor(0x26, 0x20, 0x00)
NEW        = RGBColor(0xFF, 0x80, 0x00)   # orange — build from scratch
NEW_BG     = RGBColor(0x28, 0x14, 0x00)
NEUTRAL    = RGBColor(0x44, 0x88, 0xBB)   # blue-gray — external system

SF_COL     = BLUE
SNOW_COL   = RGBColor(0x29, 0xB5, 0xE8)
ENG_COL    = RGBColor(0x99, 0x44, 0xCC)
NS_COL     = RGBColor(0x2E, 0x8B, 0x57)
RSC_COL    = RGBColor(0xE8, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 7

# ── Primitives ─────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=13, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def bg(slide, c=NAVY):   rect(slide, 0, 0, 13.33, 7.5, fill=c)

def hdr(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.0, fill=MID_DARK)
    rect(slide, 0, 1.0, 13.33, 0.06, fill=BLUE)
    txt(slide, title, 0.4, 0.08, 12.5, 0.62, size=27, bold=True)
    if sub:
        txt(slide, sub, 0.4, 0.7, 12.5, 0.32, size=12, color=TEAL)

def pgn(slide, n):
    txt(slide, f"{n} / {TOTAL}", 12.3, 7.12, 0.95, 0.3,
        size=10, color=DK_GRAY, align=PP_ALIGN.RIGHT)

def legend(slide, x, y):
    for i, (c, lbl) in enumerate([
        (BUILT,  "Already Built"),
        (EXTEND, "Extend / Modify"),
        (NEW,    "Build from Scratch"),
    ]):
        xi = x + i * 3.1
        rect(slide, xi, y, 0.22, 0.22, fill=c)
        txt(slide, lbl, xi + 0.28, y + 0.01, 2.6, 0.22, size=12, color=LT_GRAY)

# ── Flow node ─────────────────────────────────────────────────

def node(slide, x, y, w, h, label, sub=None, status="built", ls=12):
    if   status == "built":   bg_c, st = BUILT_BG,  BUILT
    elif status == "extend":  bg_c, st = EXTEND_BG, EXTEND
    elif status == "new":     bg_c, st = NEW_BG,    NEW
    else:                     bg_c, st = MID_GRAY,  NEUTRAL

    rect(slide, x, y, w, h, fill=bg_c)
    rect(slide, x, y, w, 0.07, fill=st)

    if sub:
        txt(slide, label, x+0.07, y+0.1, w-0.14, (h-0.1)*0.52,
            size=ls, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, sub, x+0.07, y+0.1+(h-0.1)*0.52, w-0.14, (h-0.1)*0.48,
            size=ls-2, italic=True, color=LT_GRAY, align=PP_ALIGN.CENTER)
    else:
        txt(slide, label, x+0.07, y+0.1, w-0.14, h-0.14,
            size=ls, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def arrow_h(slide, x, y, lbl=None, color=TEAL):
    txt(slide, "→", x, y-0.12, 0.38, 0.46, size=22, bold=True, color=color,
        align=PP_ALIGN.CENTER)
    if lbl:
        txt(slide, lbl, x-0.05, y+0.36, 0.48, 0.28, size=8, italic=True,
            color=DK_GRAY, align=PP_ALIGN.CENTER)

def arrow_v(slide, x, y, lbl=None, color=TEAL):
    txt(slide, "↓", x, y, 0.38, 0.36, size=20, bold=True, color=color,
        align=PP_ALIGN.CENTER)
    if lbl:
        txt(slide, lbl, x+0.34, y+0.04, 1.6, 0.28, size=8, italic=True, color=DK_GRAY)

def sys_box(slide, x, y, w, h, label, sub, color):
    rect(slide, x, y, w, h, fill=MID_DARK)
    rect(slide, x, y, w, 0.08, fill=color)
    txt(slide, label, x+0.1, y+0.1, w-0.2, 0.52,
        size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, sub, x+0.08, y+0.65, w-0.16, h-0.72,
        size=11, color=LT_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  S1 — TITLE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 5.8, 7.5, fill=RGBColor(0x07, 0x12, 0x1E))
rect(s, 5.8, 0, 0.07, 7.5, fill=BLUE)
rect(s, 0, 0, 5.8, 0.07, fill=TEAL)

txt(s, "RUBRIK", 0.5, 0.3, 5, 0.6, size=20, bold=True, color=TEAL)
txt(s, "MSP PayGo\nBilling Architecture", 0.5, 1.1, 5.1, 1.8,
    size=34, bold=True, color=WHITE)
txt(s, "What's Already Built\nvs.\nWhat Needs to Be Built", 0.5, 3.0, 5.1, 1.5,
    size=22, color=TEAL)
txt(s, "Visual guide · April 2026", 0.5, 6.8, 5.1, 0.5, size=13, color=DK_GRAY)

# right side — legend / key
txt(s, "Color Key", 6.3, 1.4, 6.6, 0.5, size=16, bold=True, color=WHITE)
for i, (c, lbl, desc) in enumerate([
    (BUILT,  "Green",  "Already exists in production"),
    (EXTEND, "Yellow", "Exists — needs extension/modification"),
    (NEW,    "Orange", "Needs to be built from scratch"),
    (NEUTRAL,"Blue",   "External system (Snowflake / Engineering)"),
]):
    yy = 2.05 + i * 1.2
    rect(s, 6.3, yy, 0.5, 0.65, fill=c)
    txt(s, lbl,  6.9, yy + 0.02, 5.9, 0.3, size=15, bold=True,  color=WHITE)
    txt(s, desc, 6.9, yy + 0.34, 5.9, 0.3, size=13, color=LT_GRAY)


# ══════════════════════════════════════════════════════════════
#  S2 — SYSTEM LANDSCAPE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "The Four Systems That Make PayGo Work",
           "Data flows left to right — each system has a distinct role")
pgn(s, 2)

systems = [
    (0.3,  "RSC\n(Rubrik Security Cloud)",
     "The actual product MSPs\nuse. Every backup, scan, and\nprotection event is logged.",
     RSC_COL),
    (3.55, "Snowflake\n(Data Warehouse)",
     "Stores all RSC usage logs.\nData engineers write SQL to\naggregate and calculate bills.",
     SNOW_COL),
    (6.8,  "Salesforce\n(CRM + Billing CPQ)",
     "Manages quotes, orders,\ncontracts, entitlements, and\noverage opportunity records.",
     SF_COL),
    (10.05,"NetSuite\n(Finance / ERP)",
     "Generates and sends\nthe actual invoice to the\nMSP every month.",
     NS_COL),
]
for x, label, sub, col in systems:
    sys_box(s, x, 1.3, 3.0, 4.2, label, sub, col)

# arrows between systems
for ax, lbl in [(3.35, "Usage data\nlogs"), (6.6, "Pre-calculated\noverage data\n(via Mulesoft)"), (9.85, "Order data\n(auto-sync)")]:
    arrow_h(s, ax, 3.0, lbl)

# Mulesoft callout
rect(s, 4.0, 5.85, 5.3, 1.25, fill=MID_DARK)
rect(s, 4.0, 5.85, 5.3, 0.06, fill=TEAL)
txt(s, "Mulesoft — The Messenger", 4.15, 5.93, 5.0, 0.38,
    size=13, bold=True, color=TEAL)
txt(s, "Mulesoft is middleware that moves pre-computed data from Snowflake into Salesforce via Platform Events. "
       "It also triggers the monthly billing batch job on the 26th of each month.",
    4.15, 6.35, 5.0, 0.7, size=11, color=LT_GRAY, wrap=True)

# Pacman callout
rect(s, 9.6, 5.85, 3.5, 1.25, fill=MID_DARK)
rect(s, 9.6, 5.85, 3.5, 0.06, fill=ENG_COL)
txt(s, "Pacman (RSC Provisioning)", 9.75, 5.93, 3.2, 0.38,
    size=13, bold=True, color=ENG_COL)
txt(s, "Engineering-owned system. When Salesforce says 'Order Accepted', "
       "Pacman sets up the MSP's RSC environment.",
    9.75, 6.35, 3.2, 0.7, size=11, color=LT_GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════
#  S3 — ALREADY BUILT: SCALE UTILITY PIPELINE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Already Built: The Monthly Billing Pipeline",
       "Scale Utility runs this exact cycle in production today. PayGo reuses the same pattern.")
pgn(s, 3)

txt(s, "Every month on the 26th — FULLY AUTOMATED today for Scale Utility",
    0.4, 1.1, 12.5, 0.38, size=13, bold=True, color=BUILT)

# 6-node horizontal pipeline
# 6 × 1.92" + 5 × 0.26" = 11.52 + 1.3 = 12.82 → margin 0.255"
NW, NH = 1.92, 1.35
nodes6 = [
    ("Mulesoft\ntrigger", "26th of\nevery month"),
    ("Platform Event", "Scale_Utility_\nIntegration__e"),
    ("SFDC Batch", "ScaleUtilityOverage\nBatchWithCLI.cls"),
    ("Snowflake\nFeed read", "Snowflake_Feed__c\nin Salesforce"),
    ("Overage\nRecords", "Scale_Utility_\nOverage__c"),
    ("Opp + Order\nAuto-created", "OverageOpportunity\nBatch.cls"),
]
for i, (lbl, sub) in enumerate(nodes6):
    xo = 0.26 + i * (NW + 0.26)
    node(s, xo, 1.65, NW, NH, lbl, sub, "built", ls=11)
    if i < 5:
        arrow_h(s, xo + NW + 0.03, 2.18)

# second row — result
rect(s, 0.26, 3.2, 12.82, 0.06, fill=TEAL)
txt(s, "↓   Result lands here   ↓", 5.5, 3.28, 2.8, 0.36,
    size=12, color=TEAL, align=PP_ALIGN.CENTER)

# Payment + NetSuite row
nodes2 = [
    ("Payment Schedule", "Payment_Schedule__c\nwith NS sync fields"),
    ("NetSuite Sync Batch", "Batch_SyncQueued\nTaskToNetsuite.cls"),
    ("NetSuite\nBilling Schedule", "Auto-generated\nin NetSuite"),
    ("Invoice sent\nto MSP", "NetSuite emails\ninvoice"),
]
for i, (lbl, sub) in enumerate(nodes2):
    xo = 0.5 + i * 3.1
    node(s, xo, 3.72, 2.75, NH, lbl, sub, "built", ls=11)
    if i < 3:
        arrow_h(s, xo + 2.75 + 0.04, 4.28)

# bottom note
rect(s, 0.26, 5.35, 12.82, 1.85, fill=MID_DARK)
rect(s, 0.26, 5.35, 12.82, 0.06, fill=BUILT)
txt(s, "KEY FINDING:", 0.5, 5.44, 2.0, 0.38, size=12, bold=True, color=BUILT)
txt(s, "OverageOpportunityBatch.cls already reads CPQ_MSP_Overage__c and automatically creates Closed-Won Opportunities + Orders. "
       "This covers ~80% of what PayGo needs for its monthly overage billing cycle. "
       "The team extends this — they do not rebuild it.",
    2.55, 5.44, 10.3, 0.8, size=12, color=WHITE, wrap=True)
txt(s, "CPQ_MSP_Overage__c  ·  Payment_Schedule__c (with NS fields)  ·  Batch_SyncQueuedTaskToNetsuite.cls  —  all already live in the org",
    0.5, 6.32, 12.5, 0.75, size=11, italic=True, color=LT_GRAY)


# ══════════════════════════════════════════════════════════════
#  S4 — SALESFORCE: BUILT vs TO BUILD
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Salesforce: Stage-by-Stage — Built vs. To Build",
       "The quote-to-cash pipeline inside Salesforce. Color = what work is required.")
pgn(s, 4)
legend(s, 0.4, 6.95)

# Row 1: Quote → Order stages  (top pipeline)
row1 = [
    ("Opportunity\n& Account",   None,
     "SFDC standard\nAccount + Opp objects",  "built"),
    ("CPQ Quote\nBasic",         "existing CPQ\ninfrastructure",
     None,                                     "built"),
    ("Rate Card\nSetup in CPQ",  "~21 SKUs,\nA-D categories",
     "New CPQ Product\nobject + Price Rules",  "new"),
    ("Ramp Quoting\n(PayGo SKUs)","RS-PAYGO-C\n& RS-PAYGO-NC",
     "New line items\n& quote behavior",       "new"),
    ("3-Level\nApproval",        "by commitment\ntier",
     "New SFDC\nApproval Process",             "new"),
    ("Quote PDF\nwith Rate Card", "rate card table\n+ legal terms",
     "New CPQ\ndoc template",                  "new"),
]

NW2, NH2 = 1.96, 1.45
for i, (lbl, sub, sub2, status) in enumerate(row1):
    xo = 0.28 + i * (NW2 + 0.22)
    disp_sub = sub if sub else sub2
    node(s, xo, 1.18, NW2, NH2, lbl, disp_sub, status, ls=11)
    if i < 5:
        arrow_h(s, xo + NW2 + 0.03, 1.76)

# divider
rect(s, 0.28, 2.78, 12.82, 0.05, fill=DK_GRAY)
txt(s, "Order  →  Entitlement  →  Billing", 5.3, 2.86, 3.2, 0.32,
    size=11, bold=True, color=DK_GRAY, align=PP_ALIGN.CENTER)

# Row 2: Order → Entitlement → Billing  (bottom pipeline)
row2 = [
    ("Order\nCreated",          "single-line\nPayGo order",
     "OrderHandler.cls\nextended",             "extend"),
    ("EDI / Glynt\nPO Check",   "exclude rate card\nproducts from PO",
     "Add Is_Rate_Card_\nProduct__c flag",     "extend"),
    ("Auto\nEntitlements",      "qty 99,999\nfor all 21 products",
     "EntitlementTrigger\n— new automation",   "new"),
    ("Pacman\nCallout",         "provision RSC\non Order Accepted",
     "New HTTP callout\nApex class",           "new"),
    ("Overage\nRecords",        "CPQ_MSP_Overage__c\nextended for PayGo",
     "PayGo batch\n(extends SU pattern)",      "extend"),
    ("Opp + Order\nAuto-created","Closed-Won GC\nOnDemand",
     "Extend Overage\nOpportunityBatch",       "extend"),
]

for i, (lbl, sub, sub2, status) in enumerate(row2):
    xo = 0.28 + i * (NW2 + 0.22)
    disp_sub = sub if sub else sub2
    node(s, xo, 3.25, NW2, NH2, lbl, disp_sub, status, ls=11)
    if i < 5:
        arrow_h(s, xo + NW2 + 0.03, 3.83)

# Connector arrow from row1 end to row2 start
txt(s, "↓", 0.92, 2.72, 0.4, 0.38, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  S5 — OUTSIDE SALESFORCE: SNOWFLAKE + ENGINEERING
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Outside Salesforce: Snowflake & Engineering — Built vs. To Build",
       "The computation, provisioning, and data pipelines that live outside Salesforce.")
pgn(s, 5)
legend(s, 0.4, 6.95)

# ── LEFT: Snowflake section ────────────────────────────────────
rect(s, 0.28, 1.15, 6.35, 5.55, fill=MID_DARK)
rect(s, 0.28, 1.15, 6.35, 0.07, fill=SNOW_COL)
txt(s, "Snowflake / Data Engineering", 0.42, 1.2, 6.1, 0.44, size=15, bold=True, color=SNOW_COL)

snow_nodes = [
    ("RSC Usage\nLogs arrive", "Raw data from\nproduct telemetry", "1.78", "built"),
    ("Scale Utility\nCalc (existing)", "Existing aggregation\nfor SU billing", "2.82", "built"),
    ("PayGo\nConsumption Calc", "New: aggregate by MSP,\nexclude POC, multiply\nby rate card price", "3.86", "new"),
    ("Compare vs\nCommitment", "New: overage amount\n= usage − commit\n(or $0 if no-commit)", "4.9",  "new"),
    ("Mulesoft push\nto Salesforce", "Send pre-calculated\noverage to SFDC\nSnowflake_Feed__c", "5.94", "new"),
]

NS_W, NS_H = 2.5, 0.88
for i, (lbl, sub, ystr, status) in enumerate(snow_nodes):
    yy = float(ystr)
    node(s, 0.5, yy, NS_W, NS_H, lbl, sub, status, ls=10)
    if i < 4:
        arrow_v(s, 0.93, yy + NS_H + 0.01)

# schema changes (right side of Snowflake box)
txt(s, "Schema changes needed:", 3.2, 1.75, 3.2, 0.35, size=11, bold=True, color=EXTEND)
schema_items = [
    ("Add contract_type column", "new"),
    ("Add marketing_name column", "new"),
    ("Entitlement sync: SFDC→Snowflake→RSC", "new"),
    ("Monthly billing export file", "new"),
    ("Distributor usage reports (pre-invoice)", "new"),
]
for i, (lbl, status) in enumerate(schema_items):
    yy = 2.15 + i * 0.88
    node(s, 3.15, yy, 3.25, 0.72, lbl, None, status, ls=11)

# ── RIGHT: Engineering section ────────────────────────────────
rect(s, 6.85, 1.15, 6.2, 5.55, fill=MID_DARK)
rect(s, 6.85, 1.15, 6.2, 0.07, fill=ENG_COL)
txt(s, "Engineering (RSC / Pacman)", 7.0, 1.2, 5.9, 0.44, size=15, bold=True, color=ENG_COL)

# Pacman provision flow (vertical)
eng_nodes = [
    ("Salesforce\nSends Callout",   "POST to Pacman\n/provision_order API\n(SFDC writes this)", "1.75", "new"),
    ("Pacman\nReceives Request",    "Parses order JSON\nValidates MSP account\n(Engineering builds)", "2.9", "new"),
    ("RSC Account\nProvisioned",    "MSP can now use all\n21 rate card products\nin their RSC dashboard", "4.05", "new"),
    ("Entitlement\nSync to RSC",    "SFDC entitlements →\nSnowflake → RSC\nfeature flags updated", "5.2",  "new"),
]

EN_W, EN_H = 2.6, 0.95
for i, (lbl, sub, ystr, status) in enumerate(eng_nodes):
    yy = float(ystr)
    node(s, 7.05, yy, EN_W, EN_H, lbl, sub, status, ls=10)
    if i < 3:
        arrow_v(s, 8.18, yy + EN_H + 0.01)

# NetSuite block
txt(s, "NetSuite (Finance)", 9.9, 1.75, 2.95, 0.38, size=14, bold=True, color=NS_COL)
ns_nodes = [
    ("Auto Billing\nSchedule", "On SO fulfillment\n(NetSuite config)", "built"),
    ("Consolidated\nInvoice",  "Monthly commit +\noverage merged", "new"),
    ("Usage Detail\non Invoice","Line-level breakdown\nfor MSP transparency", "new"),
]
for i, (lbl, sub, status) in enumerate(ns_nodes):
    yy = 2.2 + i * 1.3
    node(s, 9.9, yy, 2.95, 1.1, lbl, sub, status, ls=11)
    if i < 2:
        arrow_v(s, 11.27, yy + 1.1 + 0.02)


# ══════════════════════════════════════════════════════════════
#  S6 — END-TO-END FUTURE STATE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "End-to-End PayGo Flow — Future State",
       "The complete automated cycle once everything is built and live.")
pgn(s, 6)

# ── Top row: Deal sign-up track ────────────────────────────────
txt(s, "DEAL TRACK  (happens once per MSP)", 0.4, 1.15, 6, 0.3, size=10, bold=True, color=TEAL)

deal_nodes = [
    ("MSP\nsigns deal", None, "new"),
    ("CPQ Quote\n+ Rate Card", "PayGo SKU,\nramp, commit", "new"),
    ("Approval\n(1/2/3 level)", "auto-routed\nby tier", "new"),
    ("Order\nCreated", "single-line\nPayGo order", "extend"),
    ("Entitlements\nauto-created", "21 products,\nqty 99,999", "new"),
    ("Pacman\nprovisions RSC", "MSP goes live\nin RSC", "new"),
]
DW, DH = 1.82, 1.0
for i, (lbl, sub, st) in enumerate(deal_nodes):
    xo = 0.3 + i * (DW + 0.22)
    node(s, xo, 1.52, DW, DH, lbl, sub, st, ls=10)
    if i < 5:
        arrow_h(s, xo + DW + 0.03, 1.92)

# ── Divider ────────────────────────────────────────────────────
rect(s, 0.3, 2.72, 12.82, 0.05, fill=DK_GRAY)
txt(s, "MONTHLY BILLING TRACK  (automated every 26th)", 0.4, 2.8, 6, 0.3,
    size=10, bold=True, color=TEAL)

# ── Bottom row: Monthly billing track ─────────────────────────
billing_nodes = [
    ("MSP uses\nRSC product",    "usage logged\nautomatically",    "built"),
    ("Snowflake\ncalculates",    "usage × rate\nvs commitment",    "new"),
    ("Mulesoft\npushes data",    "to Snowflake_\nFeed__c in SFDC", "built"),
    ("SFDC Batch\ncreates Opp",  "Closed-Won GC\nOnDemand auto",   "extend"),
    ("NetSuite\ninvoices MSP",   "commit + overage\nmonthly bill",  "extend"),
    ("MSP pays\n& continues",    "cycle repeats\nnext month",      "neutral"),
]
BW, BH = 1.82, 1.0
for i, (lbl, sub, st) in enumerate(billing_nodes):
    xo = 0.3 + i * (BW + 0.22)
    node(s, xo, 3.18, BW, BH, lbl, sub, st, ls=10)
    if i < 5:
        arrow_h(s, xo + BW + 0.03, 3.58)

# ── Approval & exception track ─────────────────────────────────
rect(s, 0.3, 4.4, 12.82, 0.04, fill=DK_GRAY)
txt(s, "EXCEPTION TRACK  (when overage needs review before billing)", 0.4, 4.46, 8, 0.3,
    size=10, bold=True, color=EXTEND)

exc_nodes = [
    ("Overage\nexceeds threshold", None, "extend"),
    ("SFDC Approval\nProcess fires",  "manager reviews\noverage record", "new"),
    ("Approved →\nOpp auto-created", "same Closed-Won\nOpp + Order flow", "extend"),
    ("Rejected →\nReview manually",  "CSM contacts\nMSP for dispute", "neutral"),
]
EW, EH = 2.75, 0.95
for i, (lbl, sub, st) in enumerate(exc_nodes):
    xo = 0.4 + i * (EW + 0.28)
    node(s, xo, 4.82, EW, EH, lbl, sub, st, ls=11)
    if i < 3:
        arrow_h(s, xo + EW + 0.04, 5.2)

# Bottom summary strip
rect(s, 0.3, 6.0, 12.82, 1.25, fill=MID_DARK)
rect(s, 0.3, 6.0, 12.82, 0.06, fill=BLUE)
txt(s, "The entire monthly billing cycle is zero-touch once built: "
       "RSC logs usage → Snowflake calculates → Mulesoft pushes to Salesforce → "
       "Salesforce creates the Opportunity and Order → NetSuite sends the invoice. "
       "No manual intervention unless an overage dispute is raised.",
    0.5, 6.1, 12.5, 1.1, size=13, color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════════
#  S7 — SUMMARY: BUILT vs TO BUILD ACROSS ALL SYSTEMS
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Summary: What's Built vs. What Needs to Be Built",
       "Across all four systems — by team.")
pgn(s, 7)
legend(s, 0.4, 6.95)

systems_summary = [
    (SF_COL,   "Salesforce",
     ["CPQ Quote baseline", "Order handlers (OrderHandler.cls)", "Basic Entitlement trigger",
      "OverageOpportunityBatch.cls (~80% of billing)", "CPQ_MSP_Overage__c object",
      "Payment_Schedule__c → NetSuite sync", "EDI/Glynt integration base"],
     ["Rate Card object (versioned, ~21 SKUs)", "PayGo SKUs in CPQ (Commit + No-Commit)",
      "3-level approval routing", "Quote PDF with rate card", "Auto-entitlements for 21 products",
      "Pacman HTTP callout on Order Accepted", "PayGo overage batch extension"]),
    (SNOW_COL, "Snowflake / Data Eng",
     ["Raw RSC usage data storage", "Scale Utility aggregation (existing)", "Snowflake_Feed__c schema (partial)"],
     ["PayGo consumption calculation engine", "Add contract_type & marketing_name columns",
      "Entitlement sync: SFDC → Snowflake → RSC", "Monthly billing export file",
      "Distributor pre-invoice usage reports"]),
    (ENG_COL,  "Engineering (RSC/Pacman)",
     ["Pacman base provisioning system", "RSC entitlement feature system"],
     ["Pacman /provision_order API endpoint", "RSC reads marketing_name per entitlement",
      "RSC reads contract_type per entitlement", "License termination for No-Commit non-pay"]),
    (NS_COL,   "NetSuite",
     ["Basic billing schedule", "Standard invoice generation", "SFDC order data sync (via Payment_Schedule__c)"],
     ["Auto-billing schedule on SO fulfillment", "Consolidated invoice (commit + overage)",
      "Usage detail line items on invoice"]),
]

COL_W = 3.2
for ci, (col, sys_name, built_items, new_items) in enumerate(systems_summary):
    x = 0.28 + ci * (COL_W + 0.12)
    # header
    rect(s, x, 1.15, COL_W, 0.46, fill=MID_DARK)
    rect(s, x, 1.15, COL_W, 0.07, fill=col)
    txt(s, sys_name, x+0.1, 1.2, COL_W-0.2, 0.38, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Built section
    rect(s, x, 1.65, COL_W, 0.3, fill=BUILT_BG)
    rect(s, x, 1.65, COL_W, 0.05, fill=BUILT)
    txt(s, "ALREADY BUILT", x+0.08, 1.68, COL_W-0.16, 0.25, size=9, bold=True, color=BUILT)
    for bi, item in enumerate(built_items):
        rect(s, x+0.08, 2.0 + bi*0.38, 0.14, 0.14, fill=BUILT)
        txt(s, item, x+0.26, 1.98 + bi*0.38, COL_W-0.34, 0.32, size=9, color=WHITE)

    # To Build section
    sep_y = 2.0 + len(built_items)*0.38 + 0.05
    rect(s, x, sep_y, COL_W, 0.3, fill=NEW_BG)
    rect(s, x, sep_y, COL_W, 0.05, fill=NEW)
    txt(s, "NEEDS TO BE BUILT", x+0.08, sep_y+0.04, COL_W-0.16, 0.25, size=9, bold=True, color=NEW)
    for ni, item in enumerate(new_items):
        rect(s, x+0.08, sep_y+0.35 + ni*0.38, 0.14, 0.14, fill=NEW)
        txt(s, item, x+0.26, sep_y+0.33 + ni*0.38, COL_W-0.34, 0.32, size=9, color=WHITE)


# ── Save ───────────────────────────────────────────────────────
out = "/Users/subhajitbiswas/Cloud Project/RubrikClaudePOC/Consumption Based Model/PayGo_Architecture_BuiltVsToBuild.pptx"
prs.save(out)
print(f"Saved → {out}")
