from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color palette ──────────────────────────────────────────────
RUBRIK_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
RUBRIK_BLUE   = RGBColor(0x00, 0x72, 0xC6)   # bright blue
RUBRIK_TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # accent teal
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GRAY     = RGBColor(0x44, 0x55, 0x66)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank layout

# ── helpers ────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def add_bullet_box(slide, bullets, x, y, w, h, font_size=16, title=None,
                   title_color=RUBRIK_TEAL, bullet_color=WHITE, bg=None,
                   indent=0.25):
    if bg:
        add_rect(slide, x, y, w, h, fill_color=bg)
    if title:
        add_text(slide, title, x + 0.15, y + 0.1, w - 0.3, 0.4,
                 font_size=font_size + 2, bold=True, color=title_color)
        y_start = y + 0.55
    else:
        y_start = y + 0.15
    for i, b in enumerate(bullets):
        add_text(slide, f"  •  {b}", x + indent, y_start + i * (font_size / 72 + 0.28),
                 w - indent - 0.1, 0.45, font_size=font_size, color=bullet_color)

def slide_bg(slide, color=RUBRIK_DARK):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=color)

def accent_bar(slide, y=1.05, h=0.06, color=RUBRIK_BLUE):
    add_rect(slide, 0, y, 13.33, h, fill_color=color)

def section_label(slide, text, x=0.4, y=0.15):
    add_text(slide, text.upper(), x, y, 6, 0.35, font_size=11,
             bold=True, color=RUBRIK_TEAL)

def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 1, 0.35,
             font_size=10, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

TOTAL = 8

# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)

# left dark panel
add_rect(s, 0, 0, 6.5, 7.5, fill_color=RGBColor(0x06, 0x0F, 0x1A))

# blue vertical accent strip
add_rect(s, 6.5, 0, 0.08, 7.5, fill_color=RUBRIK_BLUE)

# teal horizontal strip at top
add_rect(s, 0, 0, 6.5, 0.08, fill_color=RUBRIK_TEAL)

# RUBRIK wordmark area
add_text(s, "RUBRIK", 0.5, 0.35, 5.5, 0.7,
         font_size=22, bold=True, color=RUBRIK_TEAL)

# main title
add_text(s, "MSP PayGo\nBilling Transformation", 0.5, 1.4, 5.7, 1.6,
         font_size=36, bold=True, color=WHITE)

# subtitle
add_text(s, "Run Phase — Architecture & Delivery Overview", 0.5, 3.1, 5.7, 0.6,
         font_size=16, color=LIGHT_GRAY)

# presenter line
add_text(s, "April 2026  |  Presented to CIO", 0.5, 6.8, 5.7, 0.5,
         font_size=13, color=DARK_GRAY)

# right panel — key numbers
stats = [
    ("70 → 170+", "Revised story point estimate"),
    ("4", "Engineering teams involved"),
    ("21", "Rate Card SKUs to configure"),
    ("~26th", "Monthly automated billing cycle"),
]
for i, (num, label) in enumerate(stats):
    yy = 1.5 + i * 1.4
    add_text(s, num, 7.0, yy, 5.8, 0.65,
             font_size=34, bold=True, color=RUBRIK_TEAL, align=PP_ALIGN.CENTER)
    add_text(s, label, 7.0, yy + 0.62, 5.8, 0.4,
             font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < len(stats) - 1:
        add_rect(s, 8.5, yy + 1.15, 2.8, 0.015, fill_color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — WHAT IS PAYGO?
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Business Context")
slide_number(s, 2, TOTAL)

add_text(s, "What Is PayGo?", 0.4, 0.25, 10, 0.65,
         font_size=30, bold=True, color=WHITE)

# two model cards
for i, (title, col, points) in enumerate([
    ("PayGo Commit  (RS-PAYGO-C)", RUBRIK_BLUE, [
        "MSP commits to a minimum monthly spend",
        "Pays actual usage if it exceeds the commitment",
        "Provides revenue predictability for Rubrik",
        "Unit price = $1, quantity = total commit amount",
    ]),
    ("PayGo No-Commit  (RS-PAYGO-NC)", RUBRIK_TEAL, [
        "No minimum spend guaranteed",
        "MSP pays exactly what they consume each month",
        "Zero reserve — pure consumption billing",
        "Higher flexibility, lower predictability",
    ]),
]):
    xo = 0.4 + i * 6.5
    add_rect(s, xo, 1.25, 6.1, 4.0, fill_color=RGBColor(0x10, 0x25, 0x3A))
    add_rect(s, xo, 1.25, 6.1, 0.08, fill_color=col)
    add_text(s, title, xo + 0.2, 1.4, 5.7, 0.55,
             font_size=17, bold=True, color=col)
    for j, pt in enumerate(points):
        add_text(s, f"  •  {pt}", xo + 0.2, 2.05 + j * 0.72, 5.6, 0.6,
                 font_size=14, color=WHITE)

# bottom bar
add_rect(s, 0, 5.55, 13.33, 1.65, fill_color=RGBColor(0x06, 0x0F, 0x1A))
add_text(s,
    "In both models the bill is driven by a Rate Card — a negotiated price list of ~21 products "
    "(cloud backup, ransomware protection, archival, etc.) organised into discount Categories A / B / C / D.",
    0.5, 5.65, 12.3, 1.4, font_size=15, color=LIGHT_GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — HOW SYSTEMS CONNECT TODAY
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Existing Architecture")
slide_number(s, 3, TOTAL)

add_text(s, "How the Systems Connect Today", 0.4, 0.2, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

boxes = [
    (0.3,  2.2, "RSC\n(Product)",       RUBRIK_TEAL,  "Usage data\ngenerated here"),
    (3.1,  2.2, "Snowflake\n(Data Warehouse)", RUBRIK_BLUE, "Stores & aggregates\nall usage data"),
    (5.9,  2.2, "Salesforce\n(CRM / CPQ)",     RUBRIK_BLUE, "Quotes, Orders,\nContracts, Billing"),
    (8.7,  2.2, "NetSuite\n(Finance / ERP)",    ACCENT_ORANGE, "Invoices &\npayments"),
]
for x, y, label, col, sub in boxes:
    add_rect(s, x, y, 2.5, 1.8, fill_color=RGBColor(0x10, 0x25, 0x3A))
    add_rect(s, x, y, 2.5, 0.07, fill_color=col)
    add_text(s, label, x + 0.15, y + 0.15, 2.2, 0.8,
             font_size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, sub, x + 0.1, y + 1.0, 2.3, 0.65,
             font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# arrows between boxes
arrow_labels = ["Usage logs", "via Mulesoft\n(Platform Event)", "Order sync"]
for i, lbl in enumerate(arrow_labels):
    ax = 2.85 + i * 2.78
    add_rect(s, ax, 2.95, 0.2, 0.06, fill_color=RUBRIK_BLUE)
    add_text(s, "►", ax - 0.05, 2.82, 0.4, 0.3,
             font_size=18, bold=True, color=RUBRIK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, lbl, ax - 0.3, 3.12, 0.85, 0.5,
             font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)

# Mulesoft callout
add_rect(s, 4.5, 4.3, 4.3, 0.85, fill_color=RGBColor(0x10, 0x25, 0x3A))
add_rect(s, 4.5, 4.3, 4.3, 0.05, fill_color=RUBRIK_TEAL)
add_text(s, "Mulesoft = the messenger layer between Snowflake and Salesforce",
         4.6, 4.38, 4.1, 0.7, font_size=12, color=LIGHT_GRAY, wrap=True)

add_text(s,
    "A near-identical pipeline already exists for Scale Utility billing — PayGo reuses this proven pattern.",
    0.5, 5.4, 12.3, 0.8, font_size=14, color=RUBRIK_TEAL, bold=True,
    align=PP_ALIGN.CENTER, wrap=True)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — THE SCALE UTILITY BLUEPRINT
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Existing Blueprint — Scale Utility")
slide_number(s, 4, TOTAL)

add_text(s, "We Already Have a Working Blueprint", 0.4, 0.2, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

steps = [
    ("1", "26th of every month",    "Mulesoft sends a signal\nto Salesforce"),
    ("2", "Salesforce reads data",   "Pulls pre-calculated usage\nfrom Snowflake_Feed__c"),
    ("3", "Overage records created", "Salesforce computes\nwhat each MSP owes"),
    ("4", "Auto Opportunity + Order","Salesforce creates a\nClosed-Won deal automatically"),
    ("5", "NetSuite invoices",       "Order syncs to NetSuite\nwhich sends the invoice"),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.35 + i * 2.55
    add_rect(s, x, 1.25, 2.3, 3.8, fill_color=RGBColor(0x10, 0x25, 0x3A))
    # step number circle (simulated with a small rect)
    add_rect(s, x + 0.85, 1.35, 0.6, 0.6, fill_color=RUBRIK_BLUE)
    add_text(s, num, x + 0.85, 1.33, 0.6, 0.6,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.1, 2.1, 2.1, 0.65,
             font_size=13, bold=True, color=RUBRIK_TEAL, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.1, 2.8, 2.1, 1.0,
             font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", x + 2.28, 2.7, 0.35, 0.5,
                 font_size=22, bold=True, color=RUBRIK_BLUE, align=PP_ALIGN.CENTER)

add_rect(s, 0.3, 5.3, 12.7, 1.85, fill_color=RGBColor(0x06, 0x0F, 0x1A))
add_text(s, "What this means for PayGo:", 0.55, 5.38, 6, 0.4,
         font_size=14, bold=True, color=RUBRIK_TEAL)
bullets_b = [
    "OverageOpportunityBatch.cls already creates Closed-Won GC OnDemand Opps — ~80% of the overage requirement is built",
    "CPQ_MSP_Overage__c object already exists for MSP overages",
    "Payment_Schedule__c already syncs to NetSuite — the finance pipeline is established",
]
for j, b in enumerate(bullets_b):
    add_text(s, f"  •  {b}", 0.55, 5.85 + j * 0.38, 12.5, 0.36,
             font_size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — WHAT NEEDS TO BE BUILT
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Delivery Scope")
slide_number(s, 5, TOTAL)

add_text(s, "What Needs to Be Built", 0.4, 0.2, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

cols = [
    ("IN Salesforce", RUBRIK_BLUE, 0.3, [
        "Rate Card object (~21 SKUs, versioned, A-D categories)",
        "Ramp quoting — PayGo Commit & No-Commit SKUs in CPQ",
        "3-level approval routing by commitment tier",
        "Quote document template with rate card + legal terms",
        "Auto-entitlements for all 21 products on order",
        "Extend overage batch for PayGo naming & logic",
        "Callout to Pacman (provisioning) on Order Accepted",
        "NetSuite sync extension for PayGo billing data",
    ]),
    ("IN Snowflake / Data Eng", RUBRIK_TEAL, 6.85, [
        "Monthly consumption calculation (CANNOT be in Apex — governor limits)",
        "Aggregate usage by MSP, group child accounts under parent",
        "Exclude POC and test accounts from billing",
        "Multiply usage × rate card price, compare vs commitment",
        "Add contract_type & marketing_name columns to entitlement view",
        "Entitlement sync pipeline: Salesforce → Snowflake → RSC",
        "Monthly billing export file for Finance",
        "Pre-invoice usage reports to distributors",
    ]),
]
for title, col, x, items in cols:
    add_rect(s, x, 1.2, 6.3, 5.9, fill_color=RGBColor(0x10, 0x25, 0x3A))
    add_rect(s, x, 1.2, 6.3, 0.07, fill_color=col)
    add_text(s, title, x + 0.2, 1.28, 5.9, 0.5,
             font_size=16, bold=True, color=col)
    for j, item in enumerate(items):
        add_text(s, f"  •  {item}", x + 0.15, 1.88 + j * 0.63, 5.95, 0.6,
                 font_size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — WHO BUILDS WHAT
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Team Responsibilities")
slide_number(s, 6, TOTAL)

add_text(s, "Who Builds What", 0.4, 0.2, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

teams = [
    ("Salesforce Team", RUBRIK_BLUE,
     "Rate card CPQ config, quoting, approvals, order management, "
     "entitlement automation, overage batch extension, Pacman callout, NetSuite sync extension"),
    ("Data Engineering\n(Snowflake)", RUBRIK_TEAL,
     "Consumption calculation engine, Snowflake schema changes (contract_type, "
     "marketing_name), entitlement sync pipeline, billing export file, distributor reports"),
    ("Engineering\n(RSC / Pacman)", ACCENT_ORANGE,
     "Pacman receives provisioning signals from Salesforce and sets up RSC accounts; "
     "RSC reads entitlement data for feature access control"),
    ("NetSuite Admin\n(Finance Team)", RGBColor(0xA8, 0xD8, 0xEA),
     "Auto-billing schedule creation on SO fulfillment, consolidated invoice rules, "
     "usage detail fields on invoices — all NetSuite-native configuration"),
]
for i, (team, col, desc) in enumerate(teams):
    row = i // 2
    col_x = 0.3 + (i % 2) * 6.5
    yy = 1.3 + row * 2.85
    add_rect(s, col_x, yy, 6.3, 2.5, fill_color=RGBColor(0x10, 0x25, 0x3A))
    add_rect(s, col_x, yy, 6.3, 0.07, fill_color=col)
    add_text(s, team, col_x + 0.2, yy + 0.12, 5.9, 0.6,
             font_size=16, bold=True, color=col)
    add_text(s, desc, col_x + 0.2, yy + 0.75, 5.9, 1.6,
             font_size=13, color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — ESTIMATION REALITY CHECK
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Estimation Review")
slide_number(s, 7, TOTAL)

add_text(s, "Estimation: Original vs. Revised", 0.4, 0.2, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

# table header
headers = ["Workstream", "Original", "Revised", "Key Reason for Change"]
widths  = [3.8, 1.4, 1.4, 5.9]
xs = [0.3, 4.1, 5.55, 7.0]

for i, (h, w, x) in enumerate(zip(headers, widths, xs)):
    add_rect(s, x, 1.25, w, 0.45, fill_color=RUBRIK_BLUE)
    add_text(s, h, x + 0.05, 1.27, w - 0.1, 0.4,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Rate Card Configuration",    "S / 2 pts",  "L / 8 pts",  "New versioned CPQ object, 21 SKUs, discount rules"),
    ("Quoting (CPQ)",              "M / 4 pts",  "XL / 12+ pts","Ramp quoting, 3-level approvals, SSP rules, quote doc"),
    ("Rating / Consumption Calc",  "M / 4 pts",  "Snowflake",  "Cannot be in Apex — governor limits; Snowflake owns this"),
    ("Overages",                   "L / 8 pts",  "M+ / 6 pts", "OverageOpportunityBatch already exists — extend, not rebuild"),
    ("Entitlement",                "S / 2 pts",  "M / 4 pts",  "Auto-create 21 entitlements + propagation to existing customers"),
    ("Rate Card Versioning",       "M / 4 pts",  "L / 8 pts",  "Cross-object impact: quote, order, entitlement, amendments"),
    ("Amendments",                 "M / 4 pts",  "L / 8 pts",  "Rate card interplay increases complexity"),
    ("ENGG requirements (missing)","0 pts",      "20-30 pts",  "ENGG-01–12 are P0 Must-Have — completely absent from estimate"),
    ("Snowflake scope",            "8 pts",      "25-35 pts",  "Schema changes, sync pipeline, billing export all missing"),
]
row_colors = [RGBColor(0x10, 0x25, 0x3A), RGBColor(0x0D, 0x1E, 0x30)]
for ri, row in enumerate(rows):
    yy = 1.75 + ri * 0.55
    bg = row_colors[ri % 2]
    # highlight ENGG row
    if "ENGG" in row[0] or "Snowflake" in row[0]:
        bg = RGBColor(0x3A, 0x10, 0x10)
    for ci, (val, w, x) in enumerate(zip(row, widths, xs)):
        add_rect(s, x, yy, w, 0.52, fill_color=bg)
        col_c = WHITE
        if ci == 1:  col_c = LIGHT_GRAY
        if ci == 2:  col_c = RUBRIK_TEAL
        align = PP_ALIGN.CENTER if ci in (1, 2) else PP_ALIGN.LEFT
        add_text(s, val, x + 0.05, yy + 0.06, w - 0.1, 0.42,
                 font_size=12, color=col_c, align=align)

add_text(s, "Total: ~70 pts  →  137–173 pts  (2.0–2.5× original)",
         0.3, 6.85, 12.7, 0.45, font_size=15, bold=True,
         color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — RISKS & RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, y=1.05)
section_label(s, "Risks & Recommendations")
slide_number(s, 8, TOTAL)

add_text(s, "Key Risks & Recommendations", 0.4, 0.2, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

risks = [
    (ACCENT_ORANGE, "RISK — Billing Calculation in Apex",
     "The existing Scale Utility batch already runs at batch size=1 due to Salesforce CPU limits. "
     "PayGo has ~20 products × multiple URLs × hundreds of MSPs. Building the consumption calculation "
     "in Salesforce Apex WILL fail in production.   ✓ Recommendation: Snowflake owns all calculation."),
    (ACCENT_ORANGE, "RISK — Rebuilding Instead of Extending",
     "OverageOpportunityBatch.cls and CPQ_MSP_Overage__c already implement ~80% of the overage requirement. "
     "Building a parallel system wastes 4-6 weeks and creates duplicate debt.   ✓ Recommendation: Extend existing code."),
    (RGBColor(0xCC, 0x00, 0x00), "RISK — Unbudgeted Engineering Scope (ENGG-01–12)",
     "12 P0 Must-Have Engineering requirements (Pacman provisioning, Snowflake schema, RSC sync, billing export) "
     "are completely absent from the 70-point estimate. Without these, SFDC ships but the product cannot be provisioned "
     "or billed.   ✓ Recommendation: Create separate Engineering & Data Eng budget immediately."),
    (RGBColor(0xCC, 0x00, 0x00), "RISK — NetSuite Treated as SFDC Responsibility",
     "Invoice generation, billing schedules, and consolidated invoices are NetSuite-native. "
     "Payment_Schedule__c sync already exists precisely because billing lives in NetSuite. "
     "✓ Recommendation: NetSuite Admin team owns INV-01/02/03; Salesforce only sends order data."),
]
for i, (col, title, body) in enumerate(risks):
    yy = 1.25 + i * 1.5
    add_rect(s, 0.3, yy, 12.7, 1.35, fill_color=RGBColor(0x10, 0x25, 0x3A))
    add_rect(s, 0.3, yy, 0.08, 1.35, fill_color=col)
    add_text(s, title, 0.52, yy + 0.07, 12.2, 0.45,
             font_size=14, bold=True, color=col)
    add_text(s, body, 0.52, yy + 0.52, 12.15, 0.75,
             font_size=12, color=WHITE, wrap=True)


# ── Save ───────────────────────────────────────────────────────
out = "/Users/subhajitbiswas/Cloud Project/RubrikClaudePOC/Consumption Based Model/PayGo_CIO_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
